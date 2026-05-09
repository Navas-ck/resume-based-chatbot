import os
import re
import json
import uuid
import threading
import logging
from typing import List, Dict, Tuple, Optional, Any

import numpy as np
from sentence_transformers import SentenceTransformer
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline
import torch


try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PIL import Image
    import pytesseract


    if os.name == 'nt':
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
except Exception:
    Image = None
    pytesseract = None


try:
    import faiss
except Exception:
    raise ImportError("faiss-cpu not installed. Run `pip install faiss-cpu`")

ROOT = os.getcwd()
DATA_DIR = os.path.join(ROOT, "data")
UPLOADS_DIR = os.path.join(ROOT, "uploads")
FAISS_DIR = os.path.join(ROOT, "faiss_index")
CONVO_STORE = os.path.join(ROOT, "convo_store", "conversations.jsonl")

os.makedirs(DATA_DIR, exist_ok=True)
os.makedirs(UPLOADS_DIR, exist_ok=True)
os.makedirs(os.path.dirname(CONVO_STORE), exist_ok=True)

EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
GEN_MODEL = "google/flan-t5-small"
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
MAX_SNIPPET = 400
BATCH_SIZE = 64
SCORE_THRESHOLD = 0.18

GREETINGS = {"hi", "hello", "hey", "hiya", "good morning", "good afternoon", "good evening"}

RESUME_SECTIONS = {
    "name": ["name"],
    "email": ["email", "e-mail"],
    "phone": ["phone", "mobile", "telephone", "tel"],
    "title": ["title", "designation", "position", "role", "profession", "current role"],
    "summary": ["summary", "profile", "professional summary", "about me", "overview"],
    "experience": ["experience", "work experience", "employment", "professional experience"],
    "education": ["education", "academic", "qualifications"],
    "skills": ["skills", "technical skills", "core skills"],
    "projects": ["projects", "project"],
    "achievements": ["achievements", "awards", "honors"],
}

EMAIL_PATTERN = re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b")
PHONE_PATTERN = re.compile(r"(?:\+?\d{1,3}[-.\s]?)?(?:\(?\d{2,4}\)?[-.\s]?)?[\d\-\.\s]{6,}\d")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)




def normalize_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()

def split_sentences(text: str) -> List[str]:
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"\n{2,}", "\n\n", text)
    text = re.sub(r"(?<!\n)\n(?!\n)", " ", text)
    sentences = re.split(r"(?<=[.!?;])\s+|\n\n", text.strip())
    return [s.strip() for s in sentences if s and len(s) > 2]

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, chunk_overlap: int = CHUNK_OVERLAP) -> List[str]:
    text = normalize_whitespace(text)
    if len(text) <= chunk_size:
        return [text]
    step = max(1, chunk_size - chunk_overlap)
    chunks = [text[i:i+chunk_size].strip() for i in range(0, len(text), step)]
    if len(chunks) > 1 and len(chunks[-1]) < chunk_size // 3:
        chunks[-2] = chunks[-2] + ' ' + chunks[-1]
        chunks.pop()
    return chunks

def clean_ocr_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    cleaned = []
    prev = None
    for ln in lines:
        if ln == prev:
            continue
        cleaned.append(ln)
        prev = ln
    return "\n".join(cleaned)

def extract_emails(text: str) -> List[str]:
    return list(dict.fromkeys(EMAIL_PATTERN.findall(text)))

def extract_phones(text: str) -> List[str]:
    raw = PHONE_PATTERN.findall(text)
    results = []
    seen = set()
    for candidate in raw:
        cand = re.sub(r"[\s\-\.\(\)]", "", candidate)
        if len(cand) >= 8 and cand not in seen:
            seen.add(cand)
            results.append(cand)
    return results

def first_capitalized_line(text: str) -> Optional[str]:
    for ln in text.splitlines():
        ln = ln.strip()
        if not ln:
            continue
        words = ln.split()
        cap_count = sum(1 for w in words if w[:1].isupper())
        if 1 < cap_count and len(words) <= 6:
            return ln
    return None


class FaissStore:
    def __init__(self):
        os.makedirs(FAISS_DIR, exist_ok=True)
        self.texts_path = os.path.join(FAISS_DIR, "texts.json")
        self.meta_path = os.path.join(FAISS_DIR, "metas.json")
        self.index_path = os.path.join(FAISS_DIR, "index.faiss")
        self.lock = threading.RLock()
        self.embedder = SentenceTransformer(EMBED_MODEL)
        self.dimension = self.embedder.get_sentence_embedding_dimension()
        self.index = None
        self.texts: List[str] = []
        self.metas: List[Dict] = []
        self._load_index()

    def _atomic_write_json(self, path, data):
        tmp = path + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        os.replace(tmp, path)

    def _load_index(self):
        with self.lock:
            if os.path.exists(self.texts_path):
                try:
                    self.texts = json.load(open(self.texts_path, "r", encoding="utf-8"))
                except Exception:
                    self.texts = []
            if os.path.exists(self.meta_path):
                try:
                    self.metas = json.load(open(self.meta_path, "r", encoding="utf-8"))
                except Exception:
                    self.metas = []
            if os.path.exists(self.index_path):
                try:
                    self.index = faiss.read_index(self.index_path)
                except Exception:
                    logger.warning("Failed to load FAISS index, creating new")
                    self.index = faiss.IndexFlatIP(self.dimension)
            else:
                self.index = faiss.IndexFlatIP(self.dimension)

    def _persist(self):
        with self.lock:
            self._atomic_write_json(self.texts_path, self.texts)
            self._atomic_write_json(self.meta_path, self.metas)
            faiss.write_index(self.index, self.index_path)

    def embed_texts(self, texts: List[str]) -> np.ndarray:
        if not texts:
            return np.zeros((0, self.dimension), dtype="float32")
        embs = []
        for i in range(0, len(texts), BATCH_SIZE):
            batch = texts[i:i+BATCH_SIZE]
            embs.append(self.embedder.encode(batch, convert_to_numpy=True, show_progress_bar=False))
        return np.vstack(embs).astype("float32")

    def add_texts(self, texts: List[str], metas: List[Dict]):
        new_texts = []
        new_metas = []
        existing_set = set([normalize_whitespace(t).lower() for t in self.texts])
        for t, m in zip(texts, metas):
            t_norm = normalize_whitespace(t)
            if not t_norm:
                continue
            if t_norm.lower() in existing_set:
                continue
            existing_set.add(t_norm.lower())
            new_texts.append(t_norm)
            new_metas.append(m)

        if not new_texts:
            return

        embs = self.embed_texts(new_texts)
        faiss.normalize_L2(embs)
        with self.lock:
            if self.index is None:
                self.index = faiss.IndexFlatIP(embs.shape[1])
            self.index.add(embs)
            self.texts.extend(new_texts)
            self.metas.extend(new_metas)
            self._persist()

    def search(self, query: str, k: int = 6) -> List[Tuple[int, float]]:
        if self.index is None or getattr(self.index, "ntotal", 0) == 0:
            return []
        q_emb = self.embedder.encode([query], convert_to_numpy=True, show_progress_bar=False)[0].reshape(1, -1)
        faiss.normalize_L2(q_emb)
        D, I = self.index.search(q_emb, k)
        results = []
        for idx, score in zip(I[0], D[0]):
            if idx < 0:
                continue
            results.append((int(idx), float(score)))
        return results

    def clear_index(self):
        with self.lock:
            self.index = faiss.IndexFlatIP(self.dimension)
            self.texts = []
            self.metas = []
            for p in [self.texts_path, self.meta_path, self.index_path]:
                if os.path.exists(p):
                    os.remove(p)


class ResumeExtractor:
    def __init__(self):
        self.section_keywords = RESUME_SECTIONS

    def extract_sections(self, text: str) -> Dict[str, Any]:
        text = text.replace('\r\n', '\n')
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        sections: Dict[str, List[str]] = {k: [] for k in self.section_keywords.keys()}

        current = None
        for ln in lines:
            ln_low = re.sub(r'[^a-z0-9 ]+', ' ', ln.lower())
            matched = False
            for sec, kws in self.section_keywords.items():
                if any(re.search(r'\b' + re.escape(kw) + r'\b', ln_low) for kw in kws):
                    current = sec
                    matched = True
                    break
            if matched:
                continue
            if current:
                sections[current].append(ln)

        out: Dict[str, Any] = {}

        name = None
        if sections.get('name'):
            name = ' '.join(sections['name']).strip()
        else:
            name = first_capitalized_line(text) or (lines[0] if lines else None)
        if name:
            out['name'] = normalize_whitespace(name)

        emails = extract_emails(text)
        if emails:
            out['email'] = emails[0]
        phones = extract_phones(text)
        if phones:
            out['number'] = phones[0]

        if sections.get('title'):
            out['profession'] = normalize_whitespace(' '.join(sections['title']))
        else:
            for ln in lines[:8]:
                if re.search(r'\b(engineer|developer|assistant|dentist|doctor|manager|analyst|teacher)\b', ln, flags=re.I):
                    out['profession'] = ln
                    break

        summary = '\n'.join(sections.get('summary', []))
        if not summary:
            summary = normalize_whitespace(' '.join(lines[:min(6, len(lines))]))
        out['summary'] = normalize_whitespace(summary)

        skills = '\n'.join(sections.get('skills', []))
        skills = re.sub(r'\b([A-Z ]{2,50})\b(?:\s+\1\b){1,}', lambda m: m.group(1), skills)
        out['skills'] = normalize_whitespace(skills)

        out['achievements'] = normalize_whitespace('\n'.join(sections.get('achievements', [])))
        out['education'] = normalize_whitespace('\n'.join(sections.get('education', [])))

        exp_lines = sections.get('experience', [])
        experiences: List[str] = []
        if exp_lines:
            block: List[str] = []
            for ln in exp_lines:
                if re.search(r'\b(19|20)\d{2}\b', ln) and block:
                    experiences.append(' '.join(block).strip())
                    block = [ln]
                else:
                    block.append(ln)
            if block:
                experiences.append(' '.join(block).strip())
        out['experience'] = experiences

        return out


class ChatManager:
    def __init__(self):
        self.store = FaissStore()
        self.sessions: Dict[str, List[Dict]] = {}
        self.gen_pipe = None
        self.extractor = ResumeExtractor()
        self.resume_dict: Dict[str, Any] = {}
        self._load_generator()
        self._load_data_dir()

    def _load_generator(self):
        try:
            device = 0 if torch.cuda.is_available() else -1
            tokenizer = AutoTokenizer.from_pretrained(GEN_MODEL)
            model = AutoModelForSeq2SeqLM.from_pretrained(GEN_MODEL)
            self.gen_pipe = pipeline("text2text-generation", model=model, tokenizer=tokenizer, device=device)
        except Exception as e:
            logger.warning("LLM generator unavailable: %s", e)
            self.gen_pipe = None

    def _load_data_dir(self):
        for folder in (DATA_DIR, UPLOADS_DIR):
            if not os.path.exists(folder):
                continue
            for file in os.listdir(folder):
                fp = os.path.join(folder, file)
                if os.path.isfile(fp):
                    try:
                        self.add_file(fp)
                    except Exception as e:
                        logger.warning("Failed to add file %s: %s", fp, e)

    def _read_txt(self, path): return open(path, "r", encoding="utf-8", errors="ignore").read()
    def _read_pdf(self, path):
        if PdfReader is None: return ""
        out = []
        try:
            reader = PdfReader(path)
            for p in reader.pages: out.append(p.extract_text() or "")
        except Exception: return ""
        return "\n".join(out)
    def _read_docx(self, path):
        if Document is None: return ""
        try: doc = Document(path)
        except Exception: return ""
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    def _read_pptx(self, path):
        if Presentation is None: return ""
        try:
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): texts.append(shape.text.strip())
            return "\n".join(texts)
        except Exception: return ""
    def _read_image(self, path):
        if Image is None or pytesseract is None: return ""
        try: return clean_ocr_text(pytesseract.image_to_string(Image.open(path)))
        except Exception: return ""
    def _read_file(self, path):
        n = path.lower()
        if n.endswith('.txt'): return self._read_txt(path)
        if n.endswith('.pdf'): return self._read_pdf(path)
        if n.endswith('.docx'): return self._read_docx(path)
        if n.endswith('.pptx'): return self._read_pptx(path)
        if n.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')): return self._read_image(path)
        return ""

    def add_file(self, filepath: str):
        filepath = os.path.abspath(filepath)
        fname = os.path.basename(filepath)
        logger.info("Adding file: %s", fname)
        text = self._read_file(filepath)
        if not text.strip():
            logger.info("No text extracted from %s, skipping.", fname)
            return False
        text = clean_ocr_text(text)
        resume_keywords = ["resume", "cv", "curriculum vitae", "experience", "education", "skills"]
        is_resume = any(k in text[:2000].lower() for k in resume_keywords)
        if is_resume:
            sections = self.extractor.extract_sections(text)
            self.resume_dict = sections
            texts_to_add, metas_to_add = [], []
            for sec, val in sections.items():
                if isinstance(val, list):
                    for i, entry in enumerate(val):
                        if entry.strip():
                            texts_to_add.append(entry.strip())
                            metas_to_add.append({"source": fname, "section": sec, "chunk": i})
                else:
                    if str(val).strip():
                        texts_to_add.append(str(val).strip())
                        metas_to_add.append({"source": fname, "section": sec, "chunk": 0})
            self.store.add_texts(texts_to_add, metas_to_add)
        else:
            chunks = chunk_text(text)
            metas = [{"source": fname, "section": "blob", "chunk": i} for i in range(len(chunks))]
            self.store.add_texts(chunks, metas)
        return True

    def _complete_sentence(self, text: str) -> str:
        if not text: return ""
        m = re.search(r"([\s\S]*?[\.\?!])", text)
        return normalize_whitespace(m.group(1)) if m else split_sentences(text)[0] if split_sentences(text) else normalize_whitespace(text)

    def _dict_reply(self, question: str) -> Optional[str]:
        q = question.lower().strip()
        rd = self.resume_dict
        if not rd: return None

        name = rd.get('name', '').lower() if rd.get('name') else ''
        if name and name in q:
            parts = []
            if rd.get('summary'): parts.append(rd.get('summary'))
            if rd.get('experience'):
                parts.extend(rd.get('experience'))
            return self._complete_sentence(' '.join(parts))

        if any(tok in q for tok in ["profession", "what does she do", "what does he do"]):
            return self._complete_sentence(str(rd.get('profession', '')))

        if 'experience' in q and rd.get('experience'):
            return '\n'.join([self._complete_sentence(exp) for exp in rd.get('experience', [])])

        if 'skill' in q:
            sk = rd.get('skills', '')
            return self._complete_sentence(sk)
        if 'education' in q and rd.get('education'):
            return '\n'.join([self._complete_sentence(edu) for edu in rd.get('education', []).split('\n')])

        if any(tok in q for tok in ['achievement', 'award', 'honor', 'recognition']) and rd.get('achievements'):
            return '\n'.join([self._complete_sentence(a) for a in rd.get('achievements', []).split('\n')])

        if any(tok in q for tok in ["profile", "summary", "about"]):
            return self._complete_sentence(rd.get('summary', ''))

        if 'name' in q: return rd.get('name', '')
        if 'email' in q: return rd.get('email', '')
        if any(tok in q for tok in ['phone', 'mobile', 'contact']): return rd.get('number', '')

        return None

    def generate_answer(self, session_id: Optional[str], query: str, k: int = 25) -> Tuple[str, List[Dict]]:
        if not session_id: session_id = uuid.uuid4().hex
        if session_id not in self.sessions: self.sessions[session_id] = []
        q = query.strip()
        if not q: return "Please ask something.", []

        if q.lower() in GREETINGS:
            reply = "Hello! How can I help you today?"
            self.sessions[session_id].append({"role": "assistant", "text": reply})
            return reply, []

        dict_reply = self._dict_reply(q)
        if dict_reply:
            self.sessions[session_id].append({"role": "assistant", "text": dict_reply})
            return dict_reply, []

        retrieved = self.store.search(q, k)
        contexts = [{"index": idx, "score": score, "text": self.store.texts[idx], "meta": self.store.metas[idx]} for idx, score in retrieved]
        if not contexts: return "No relevant info found.", []

        contexts = [c for c in contexts if c["score"] >= SCORE_THRESHOLD] if any(c["score"] >= SCORE_THRESHOLD for c in contexts) else contexts
        contexts = sorted(contexts, key=lambda x: x["score"], reverse=True)
        ctx_texts = []
        used_sections = set()
        for c in contexts:
            sec = c["meta"].get("section", "blob")
            if sec in used_sections and len(ctx_texts) >= 6: continue
            used_sections.add(sec)
            ctx_texts.append(f"[{sec}] {c['text']}")
        context_text = "\n".join(ctx_texts)[:2500]

        long_request = bool(re.search(r"\b(explain|detailed|summarize|elaborate|describe)\b", q, flags=re.I))
        if not long_request:
            short_prompt = f"Using only the context below, answer concisely.\nContext:\n{context_text}\nQuestion: {q}\nAnswer (one sentence):"
            if self.gen_pipe:
                try: out = self.gen_pipe(short_prompt, max_new_tokens=80); reply = out[0]["generated_text"].strip()
                except Exception: reply = contexts[0]["text"][:MAX_SNIPPET]
            else: reply = split_sentences(contexts[0]["text"])[0][:MAX_SNIPPET]
            reply = self._complete_sentence(reply)
            self.sessions[session_id].append({"role": "assistant", "text": reply})
            return reply, contexts

        long_prompt = f"Using only the context below, provide multi-line answer with sections.\nContext:\n{context_text}\nQuestion: {q}\nAnswer:"
        if self.gen_pipe:
            try: out = self.gen_pipe(long_prompt, max_new_tokens=300); reply = out[0]["generated_text"].strip()
            except Exception: reply = "\n\n".join([f"[{c['meta']['section']}] {c['text'][:400]}" for c in contexts[:3]])
        else: reply = "\n\n".join([f"[{c['meta']['section']}] {c['text'][:400]}" for c in contexts[:3]])
        reply = '\n'.join([self._complete_sentence(p) for p in reply.split('\n\n') if p.strip()])
        self.sessions[session_id].append({"role": "assistant", "text": reply})
        return reply, contexts

 
    def clear_all_uploads(self):
        for folder in [UPLOADS_DIR]:
            for f in os.listdir(folder):
                path = os.path.join(folder, f)
                if os.path.isfile(path): os.remove(path)
        logger.info("All uploads cleared.")

    def clear_all_data(self):
        for folder in [DATA_DIR]:
            for f in os.listdir(folder):
                path = os.path.join(folder, f)
                if os.path.isfile(path): os.remove(path)
        logger.info("All data files cleared.")

    def clear_faiss(self):
        self.store.clear_index()
        logger.info("FAISS index cleared.")

    def clear_everything(self):
        self.clear_all_uploads()
        self.clear_all_data()
        self.clear_faiss()
        self.resume_dict = {}
        logger.info("All uploads, data, FAISS, and resume dict cleared.")

    def clear_session(self, session_id: str):
        if session_id in self.sessions:
            self.sessions[session_id] = []
            logger.info(f"Session {session_id} cleared.")

    def clear_all_sessions(self):
        self.sessions = {}
        logger.info("All sessions cleared.")


chat_manager = ChatManager()

if __name__ == '__main__':
    print("ChatManager ready. Add files to ./data or ./uploads and call chat_manager.generate_answer(session_id, question)")
